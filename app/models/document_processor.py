import os
import json
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import torch
import re
from app.database.db import add_document, add_chunks, search_similar_chunks
import numpy as np

class DocumentProcessor:
    def __init__(self):
        # Chargement du modèle d'embedding local
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
    def process_document(self, filepath, filename):
        """Traite un document et l'ajoute à la base de données"""
        # Extraction du texte selon le type de fichier
        extension = os.path.splitext(filename)[1].lower()
        
        if extension == '.pdf':
            content_text = self._extract_pdf_text(filepath)
        elif extension == '.docx':
            content_text = self._extract_docx_text(filepath)
        elif extension == '.pptx':
            content_text = self._extract_pptx_text(filepath)
        else:
            raise ValueError(f"Format de fichier non pris en charge: {extension}")
        
        # Découpage du texte en chunks
        chunks = self._split_into_chunks(content_text)
        
        # Génération des embeddings pour chaque chunk
        embeddings = self._generate_embeddings(chunks)
        
        # Stockage dans la base de données
        document_id = add_document(filename, filepath, content_text)
        add_chunks(document_id, chunks, [json.dumps(emb.tolist()) for emb in embeddings])
        
        return document_id
        
    def _extract_pdf_text(self, filepath):
        """Extrait le texte d'un fichier PDF"""
        text = ""
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n\n"
        return text
    
    def _extract_docx_text(self, filepath):
        """Extrait le texte d'un fichier DOCX"""
        doc = DocxDocument(filepath)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    
    def _extract_pptx_text(self, filepath):
        """Extrait le texte d'un fichier PPTX"""
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n\n"
        return text
    
    def _split_into_chunks(self, text, max_tokens=512):
        """Découpe intelligemment le texte en chunks en préservant la structure du document"""
        # Diviser le document en sections logiques
        # Utiliser les sauts de ligne multiples, titres, et autres marqueurs structurels
        sections = re.split(r'\n\s*\n', text)
        
        chunks = []
        current_chunk = ""
        current_length = 0
        current_heading = ""
        estimated_token_size = 4  # Approximation: 1 token ≈ 4 caractères
        
        # Identifier les titres dans le document
        is_heading = lambda s: (
            re.match(r'^#+\s+', s) or                # Markdown heading
            re.match(r'^[A-Z0-9\s]{5,40}$', s) or    # ALL CAPS or mostly caps title
            (len(s) < 80 and                         # Short line
             (s.isupper() or                         # ALL CAPS
              bool(re.match(r'^[IVX0-9]+\.', s)) or  # Roman numerals or numbers at start
              bool(re.match(r'^[A-Z][\w\s]+:$', s))  # Title with colon at end
             )
            )
        )
        
        # Analyse améliorée du contenu
        for section in sections:
            # Nettoyer la section
            section = section.strip()
            if not section:
                continue
            
            # Estimer la taille de la section en tokens
            section_size = len(section) / estimated_token_size
            
            # Détecter si c'est un titre
            section_is_heading = is_heading(section)
            
            # Préserver les titres avec leur contenu pour un meilleur contexte
            if section_is_heading:
                current_heading = section
                # Si le chunk actuel n'est pas vide et qu'on rencontre un nouveau titre,
                # finaliser le chunk actuel
                if current_length > 0:
                    chunks.append(current_chunk.strip())
                    current_chunk = current_heading + "\n\n"
                    current_length = (len(current_heading) + 2) / estimated_token_size
                else:
                    # Commencer un nouveau chunk avec le titre
                    current_chunk = current_heading + "\n\n"
                    current_length = (len(current_heading) + 2) / estimated_token_size
                continue
            
            # Si la section est trop grande pour un seul chunk
            if section_size > max_tokens:
                # Si le chunk actuel n'est pas vide, l'ajouter d'abord
                if current_length > 0:
                    # Ajouter le titre au chunk seulement s'il existe
                    chunk_to_add = current_chunk.strip()
                    if chunk_to_add:
                        chunks.append(chunk_to_add)
                    current_chunk = ""
                    current_length = 0
                
                # Diviser la section en phrases
                sentences = re.split(r'(?<=[.!?])\s+', section)
                
                # Regrouper les phrases en chunks logiques
                temp_chunk = current_heading + "\n\n" if current_heading else ""
                temp_size = (len(temp_chunk)) / estimated_token_size
                
                for sentence in sentences:
                    sentence_size = len(sentence) / estimated_token_size
                    
                    # Si la phrase est trop grande, la diviser en fragments plus petits
                    if sentence_size > max_tokens:
                        # Ajouter le chunk actuel s'il n'est pas vide
                        if temp_size > 0:
                            chunks.append(temp_chunk.strip())
                        
                        # Diviser la phrase en fragments fixes
                        parts = self._split_large_sentence(sentence, max_tokens, estimated_token_size)
                        # Ajouter chaque fragment comme chunk séparé, avec contexte du titre
                        for part in parts:
                            if current_heading:
                                chunks.append(f"{current_heading}\n\n{part}".strip())
                            else:
                                chunks.append(part.strip())
                        
                        # Réinitialiser le chunk temporaire
                        temp_chunk = current_heading + "\n\n" if current_heading else ""
                        temp_size = (len(temp_chunk)) / estimated_token_size
                    else:
                        # Ajouter la phrase au chunk actuel ou créer un nouveau chunk
                        if temp_size + sentence_size <= max_tokens:
                            temp_chunk += " " + sentence if temp_chunk.endswith("\n\n") else sentence
                            temp_size += sentence_size
                        else:
                            # Finaliser le chunk actuel
                            if temp_chunk:
                                chunks.append(temp_chunk.strip())
                            # Commencer un nouveau chunk, inclure le titre si disponible
                            temp_chunk = current_heading + "\n\n" + sentence if current_heading else sentence
                            temp_size = (len(temp_chunk)) / estimated_token_size
                
                # Ajouter le dernier chunk temporaire s'il n'est pas vide
                if temp_chunk and not temp_chunk.endswith("\n\n"):
                    chunks.append(temp_chunk.strip())
                
                # Continuer avec un nouveau chunk vide
                current_chunk = ""
                current_length = 0
            else:
                # La section entière peut tenir dans un chunk
                # Vérifier si l'ajout de cette section dépasse la taille maximale
                if current_length + section_size <= max_tokens:
                    # Si c'est le début d'un chunk ou si le chunk actuel se termine par un titre,
                    # pas besoin d'espace supplémentaire
                    if not current_chunk or current_chunk.endswith("\n\n"):
                        current_chunk += section
                    else:
                        current_chunk += "\n\n" + section
                    current_length += section_size
                else:
                    # Finaliser le chunk actuel
                    chunks.append(current_chunk.strip())
                    # Commencer un nouveau chunk avec le titre si disponible
                    current_chunk = current_heading + "\n\n" + section if current_heading else section
                    current_length = (len(current_chunk)) / estimated_token_size
        
        # Ajouter le dernier chunk s'il n'est pas vide
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Assurer que chaque chunk a un contexte suffisant
        final_chunks = []
        for i, chunk in enumerate(chunks):
            # Ajouter un identifiant de position pour référence
            final_chunks.append(chunk)
        
        return final_chunks
    
    def _split_large_sentence(self, sentence, max_tokens, estimated_token_size):
        """Divise une phrase très longue en fragments plus petits de manière intelligente"""
        parts = []
        max_chars = int(max_tokens * estimated_token_size)
        
        # Essayer de diviser aux virgules d'abord
        comma_splits = re.split(r',\s+', sentence)
        
        if len(comma_splits) > 1:
            current_part = ""
            for split in comma_splits:
                if len(current_part) + len(split) + 2 <= max_chars:  # +2 pour ", "
                    current_part += (", " + split) if current_part else split
                else:
                    if current_part:
                        parts.append(current_part)
                    current_part = split
            if current_part:
                parts.append(current_part)
        else:
            # Si pas de virgules, diviser par taille fixe
            for i in range(0, len(sentence), max_chars - 10):  # -10 pour éviter de couper au milieu d'un mot
                part = sentence[i:i + max_chars - 10]
                if i + max_chars - 10 < len(sentence):
                    # Trouver le dernier espace pour éviter de couper un mot
                    last_space = part.rfind(' ')
                    if last_space > max_chars * 0.8:  # S'assurer qu'on a une portion significative
                        part = part[:last_space]
                parts.append(part)
        
        return parts
    
    def _generate_embeddings(self, chunks):
        """Génère des embeddings pour chaque chunk de texte"""
        embeddings = self.embedding_model.encode(chunks, convert_to_tensor=True)
        return embeddings 

def search_similar_chunks(query_embedding, limit=5):
    """Recherche les chunks les plus similaires à la requête"""
    query_embedding_array = json.loads(query_embedding)
    db = get_db()
    
    cursor = db.execute(
        '''SELECT c.id, c.content, c.document_id, d.title, c.embedding
        FROM chunks c JOIN documents d ON c.document_id = d.id'''
    )
    
    chunks = cursor.fetchall()
    similarities = []
    
    for chunk in chunks:
        chunk_embedding = json.loads(chunk['embedding'])
        # Calcul de similarité cosinus
        similarity = np.dot(query_embedding_array, chunk_embedding) / (
            np.linalg.norm(query_embedding_array) * np.linalg.norm(chunk_embedding)
        )
        similarities.append((chunk, similarity))
    
    # Trier par similarité décroissante
    similarities.sort(key=lambda x: x[1], reverse=True)
    
    # Retourner les chunks les plus similaires
    return [dict(s[0]) for s in similarities[:limit]] 